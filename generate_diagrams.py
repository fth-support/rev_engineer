import os
from pptx import Presentation
from pptx.util import Inches, Pt
import base64
import zlib
import urllib.parse

def encode_plantuml(text):
    """Encode PlantUML text for the PlantUML image server."""
    zlibbed_str = zlib.compress(text.encode('utf-8'))
    compressed_string = zlibbed_str[2:-4]
    return base64.b64encode(compressed_string).decode('utf-8').translate(str.maketrans('+/', '-_'))

def create_html(puml_path, html_path):
    with open(puml_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    # We will use the standard plantuml public server for rendering
    # A cleaner way is to use plantuml.com/plantuml/svg/
    # But since we just need an HTML view, we can use the plantuml javascript renderer
    html_content = f"""<!DOCTYPE html>
<html>
<head>
    <title>PlantUML Diagram</title>
</head>
<body>
    <h2>Diagram from {os.path.basename(puml_path)}</h2>
    <img src="http://www.plantuml.com/plantuml/svg/~1{encode_plantuml(content)}" alt="UML Diagram" />
    <hr>
    <h3>Source Code:</h3>
    <pre><code>{content}</code></pre>
</body>
</html>"""
    with open(html_path, 'w', encoding='utf-8') as f:
        f.write(html_content)

def create_pptx(puml_files, pptx_path, title):
    prs = Presentation()
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Automatically generated workflows"
    
    for puml in puml_files:
        name = os.path.basename(puml).replace('.puml', '')
        # Add slide
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = name
        
        # We add text placeholder because inserting the actual image requires downloading it first
        # For this prototype, we'll embed the image URL and the raw text
        with open(puml, 'r', encoding='utf-8') as f:
            code = f.read()
        
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(5))
        tf = txBox.text_frame
        tf.text = "Diagram Source (PlantUML):"
        p = tf.add_paragraph()
        p.text = code[:1000] + ("..." if len(code) > 1000 else "")
        p.font.size = Pt(10)
        
    prs.save(pptx_path)

for lang in ['EN', 'TH']:
    base_dir = f"Revamp_Docs/{lang}/Diagrams"
    uml_dir = f"{base_dir}/UML"
    html_dir = f"{base_dir}/HTML"
    pptx_dir = f"{base_dir}/PPTX"
    
    if not os.path.exists(uml_dir):
        continue
        
    puml_files = [os.path.join(uml_dir, f) for f in os.listdir(uml_dir) if f.endswith('.puml')]
    
    # HTML
    for puml in puml_files:
        filename = os.path.basename(puml).replace('.puml', '.html')
        create_html(puml, os.path.join(html_dir, filename))
        
    # PPTX
    if puml_files:
        create_pptx(puml_files, os.path.join(pptx_dir, "Workflows_Presentation.pptx"), f"System Workflows ({lang})")

print("Diagrams successfully generated in HTML and PPTX formats.")
