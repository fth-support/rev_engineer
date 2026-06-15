# -*- coding: utf-8 -*-
"""Inject POSFront (upstream producer) context into the existing doc-claude-ver
deliverables, then refresh the website bundle + download zip.

The original docx/pptx generators were removed, so this is a *patch* over the
existing files (preserving their styling). It is idempotent: re-running first
removes any previously-injected POSFront block / slide before re-adding it.

Adds, in each artifact:
  * DOCX (x3): a "POSFront (Upstream Producer)" section after the Executive Summary
  * PPTX     : one themed slide after the Executive Summary
Then copies doc-claude-ver/{Diagrams,*.docx,*.pptx} -> knowledge-pool/public/doc-claude-ver
and rebuilds knowledge-pool/public/downloads/doc-claude-ver.zip.

All facts verified against legacy_source/POSFront (mDB.bas Sub Main,
Center/mCNSP.bas trigger TRG_Tmp2Sale -> STP_PRCxTmp2Sale, cCNDB.cls).
"""
import os, re, shutil, zipfile
import docx
from docx.shared import Pt, RGBColor, Inches
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PInches, Pt as PPt
from pptx.dml.color import RGBColor as PColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

BASE = os.path.normpath(os.path.join(os.path.dirname(__file__), '..'))
PUBLIC = os.path.normpath(os.path.join(BASE, '..', 'knowledge-pool', 'public', 'doc-claude-ver'))
DOWNLOADS = os.path.normpath(os.path.join(BASE, '..', 'knowledge-pool', 'public', 'downloads'))

MARK = '​PF'  # zero-width marker to make the injection idempotent

# --------------------------------------------------------------------------- DOCX
DOCX_FILES = ['01_SRS_Document.docx', '02_Data_Dictionary.docx', '03_Program_Specification.docx']

DOCX_SECTION = {
    'title': 'บริบทต้นทาง: POSFront (Upstream Producer)',
    'intro': ('บนเครื่อง POS มีสองโปรแกรม VB6 ทำงานร่วมกันแบบ Producer / Consumer และส่งต่องานกันผ่าน Local DB '
              'ด้วยธงสถานะ 2 ตัว — ServiceTransfer ไม่ได้สร้างยอดขายเอง แต่รับช่วงต่อจาก POSFront '
              '(โปรแกรมขายหน้าร้านของ Adasoft โครงการ TAKASHIMAYA, v6.2412.1) ทำให้เครื่อง POS ทำงาน Offline ได้เต็มรูปแบบ'),
    'orch_h': 'Orchestration — POSFront เป็นผู้เปิดบริการเบื้องหลัง',
    'orch': [
        'POSFront.Sub Main (mDB.bas) เปิดบริการ 4 ตัว: ServiceOnOff.exe (สถานะ Online/Offline)',
        'ServiceAutoReplicate.exe — Replicate Master จาก Central Server ลง Local',
        'ServiceAutoClear.exe — เคลียร์ไฟล์/ข้อมูลค้าง',
        'ServiceTranfer.exe — คือ ServiceTransfer; เป็น Child process ของ POSFront ไม่ใช่โปรแกรมเดี่ยว',
    ],
    'hand_h': 'The Two-Flag Handshake (สัญญาระหว่างสองโปรแกรม)',
    'hand': [
        "FTShdStaDoc (เจ้าของ = POSFront): '2' กำลังทำ -> '1' เสร็จสมบูรณ์ เมื่อจ่ายเงิน + พิมพ์ใบเสร็จ",
        "FTStaSentOnOff (เจ้าของ = ServiceTransfer): '0' รอส่ง -> '1' ส่งขึ้น HQ แล้ว",
        ("เมื่อ FTShdStaDoc ถูก UPDATE '2'->'1' บนตารางทำงาน TPSHD<Tmn> -> SQL Trigger TRG_Tmp2Sale<Tmn> "
         "(AFTER UPDATE) เรียก STP_PRCxTmp2Sale ย้ายบิลเข้า TPSTSalHD/DT/RC/CD (ตั้ง FTStaSentOnOff='0' รอซิงค์)"),
    ],
}


def styles_by_name(doc):
    """Built-in heading styles can fail name-lookup; resolve by iterating instead."""
    return {s.name: s for s in doc.styles if s.name}


def remove_prior_injection(doc):
    """Delete a previously-injected block (paragraphs carrying the marker)."""
    for p in list(doc.paragraphs):
        if MARK in p.text:
            p._element.getparent().remove(p._element)


def find_anchor(doc):
    """First Heading-1 whose text starts with '1' (i.e. the Introduction/section 1)."""
    for p in doc.paragraphs:
        sn = p.style.name if p.style is not None else ''
        if sn == 'Heading 1' and re.match(r'^\s*1[\.\s]', p.text):
            return p
    return None


def patch_docx(path):
    doc = docx.Document(path)
    remove_prior_injection(doc)
    anchor = find_anchor(doc)
    smap = styles_by_name(doc)
    body_style = smap.get('List Paragraph')

    def ins_para(text, style=None, bold=False, color=None, size=None):
        style_obj = smap.get(style) if style else None
        if anchor is not None:
            p = anchor.insert_paragraph_before()
        else:
            p = doc.add_paragraph()
        if style_obj is not None:
            p.style = style_obj
        run = p.add_run(text)
        run.font.name = 'Tahoma'
        if bold:
            run.font.bold = True
        if color:
            run.font.color.rgb = RGBColor.from_string(color)
        if size:
            run.font.size = Pt(size)
        # carry the idempotency marker invisibly
        m = p.add_run(MARK)
        m.font.size = Pt(1)
        return p

    ins_para(DOCX_SECTION['title'], style='Heading 1', color='0D47A1')
    ins_para(DOCX_SECTION['intro'])
    ins_para(DOCX_SECTION['orch_h'], style='Heading 2')
    for b in DOCX_SECTION['orch']:
        ins_para('•  ' + b, style='List Paragraph')
    ins_para(DOCX_SECTION['hand_h'], style='Heading 2')
    for b in DOCX_SECTION['hand']:
        ins_para('•  ' + b, style='List Paragraph')
    doc.save(path)
    print('patched', os.path.basename(path))


# --------------------------------------------------------------------------- PPTX
NAVY, BLUE, LBLUE = '0D47A1', '1565C0', 'E3F2FD'
GREEN, LGREEN = '2E7D32', 'E8F5E9'
DARK, MUTED, CARD = '1A2733', '5A6B7B', 'F4F7FA'
PURPLE = '6A1B9A'


def _txt(slide, l, t, w, h, text, size, color, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(PInches(l), PInches(t), PInches(w), PInches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = PInches(0.04)
    tf.margin_top = tf.margin_bottom = PInches(0.02)
    para = tf.paragraphs[0]
    para.alignment = align
    r = para.add_run()
    r.text = text
    r.font.name = 'Tahoma'
    r.font.size = PPt(size)
    r.font.bold = bold
    r.font.color.rgb = PColor.from_string(color)
    return tb


def _rect(slide, l, t, w, h, fill, rounded=True, line=None):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        PInches(l), PInches(t), PInches(w), PInches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = PColor.from_string(fill)
    if line:
        shp.line.color.rgb = PColor.from_string(line)
        shp.line.width = PPt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _card(slide, l, t, w, h, accent, lfill, head, sub, bullets):
    _rect(slide, l, t, w, h, lfill)
    _rect(slide, l, t, 0.09, h, accent, rounded=False)
    _txt(slide, l + 0.22, t + 0.14, w - 0.4, 0.32, head, 14, DARK, bold=True)
    _txt(slide, l + 0.22, t + 0.48, w - 0.4, 0.26, sub, 10, accent, bold=True)
    y = t + 0.82
    for b in bullets:
        _txt(slide, l + 0.22, y, w - 0.42, 0.3, '•  ' + b, 9.5, MUTED)
        y += 0.34


def remove_prior_slide(prs):
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    for idx, slide in enumerate(prs.slides):
        if any(sh.has_text_frame and MARK in sh.text_frame.text for sh in slide.shapes):
            sldIdLst.remove(ids[idx])
            return


def patch_pptx(path):
    prs = Presentation(path)
    remove_prior_slide(prs)
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    _rect(slide, 0, 0, 10, 5.625, 'FFFFFF', rounded=False)

    _txt(slide, 0.55, 0.32, 9, 0.3, 'UPSTREAM PRODUCER', 11, BLUE, bold=True)
    _txt(slide, 0.55, 0.58, 9, 0.55, 'POSFront — โปรแกรมขายต้นทาง', 25, DARK, bold=True)
    _txt(slide, 0.55, 1.28, 8.9, 0.5,
         'บนเครื่อง POS มีสองโปรแกรมทำงานร่วมกันแบบ Producer / Consumer ส่งต่อกันผ่าน Local DB ด้วยธงสถานะ 2 ตัว',
         13, NAVY, bold=True)

    _card(slide, 0.55, 1.95, 4.42, 1.95, BLUE, LBLUE, 'POSFront (Producer)', 'โปรแกรมขายหน้าร้าน · VB6 v6.2412.1',
          ['อ่าน Master จาก Server', "เขียนบิลลง TPSHD<Tmn> (StaDoc='2')",
           "จ่ายเงิน+พิมพ์ใบเสร็จ -> StaDoc='1'"])
    _card(slide, 5.03, 1.95, 4.42, 1.95, PURPLE, 'F3E5F5', 'ServiceTransfer (Consumer)', 'เอเจนต์ซิงค์ · Timer 500ms',
          ["poll FTStaSentOnOff='0'", 'Tokenize -> push ขึ้น Central DB',
           "set FTStaSentOnOff='1' (synced)"])

    _rect(slide, 0.55, 4.05, 8.9, 1.18, CARD)
    _txt(slide, 0.75, 4.16, 8.5, 0.3,
         'Orchestration:  POSFront.Sub Main เปิด ServiceOnOff · ServiceAutoReplicate · ServiceAutoClear · ServiceTranfer',
         10.5, DARK, bold=True)
    _txt(slide, 0.75, 4.52, 8.5, 0.62,
         ("Handshake:  FTShdStaDoc ('2'->'1') = POSFront  |  FTStaSentOnOff ('0'->'1') = ServiceTransfer  |  "
          'Trigger TRG_Tmp2Sale<Tmn> -> STP_PRCxTmp2Sale ย้ายบิลเข้า TPSTSal*'),
         9.5, MUTED)

    # idempotency marker (tiny, off-canvas-ish but present)
    mk = _txt(slide, 0.0, 5.55, 0.3, 0.1, MARK, 1, 'FFFFFF')

    # move the new slide to right after the Executive Summary (index 3)
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    sldIdLst.remove(ids[-1])
    sldIdLst.insert(min(3, len(ids) - 1), ids[-1])
    prs.save(path)
    print('patched', os.path.basename(path))


# --------------------------------------------------------------------------- bundle
def refresh_bundle():
    # copy Diagrams (overwrite in place — OneDrive can lock the dir against rmtree)
    src_dia = os.path.join(BASE, 'Diagrams')
    dst_dia = os.path.join(PUBLIC, 'Diagrams')
    os.makedirs(dst_dia, exist_ok=True)
    for fn in os.listdir(src_dia):
        shutil.copy2(os.path.join(src_dia, fn), os.path.join(dst_dia, fn))
    # copy docx/pptx
    for f in DOCX_FILES + ['ServiceTransfer_Complete_Presentation.pptx']:
        shutil.copy2(os.path.join(BASE, f), os.path.join(PUBLIC, f))
    print('bundle refreshed ->', PUBLIC)
    # rebuild zip
    os.makedirs(DOWNLOADS, exist_ok=True)
    zpath = os.path.join(DOWNLOADS, 'doc-claude-ver.zip')
    if os.path.exists(zpath):
        os.remove(zpath)
    with zipfile.ZipFile(zpath, 'w', zipfile.ZIP_DEFLATED) as z:
        for root, _, files in os.walk(PUBLIC):
            for fn in files:
                fp = os.path.join(root, fn)
                arc = os.path.join('doc-claude-ver', os.path.relpath(fp, PUBLIC))
                z.write(fp, arc)
    print('zip rebuilt ->', zpath)


if __name__ == '__main__':
    for f in DOCX_FILES:
        patch_docx(os.path.join(BASE, f))
    patch_pptx(os.path.join(BASE, 'ServiceTransfer_Complete_Presentation.pptx'))
    refresh_bundle()
    print('done.')
